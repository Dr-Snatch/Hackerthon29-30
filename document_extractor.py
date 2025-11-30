"""
Document Extractor - Extract text from various document formats
================================================================

Supports:
- PDF (.pdf)
- Word Documents (.doc, .docx)
- PowerPoint Presentations (.ppt, .pptx)
- Plain Text (.txt)

Note: .doc and .ppt (older binary formats) may have limited support.
For best results, use .docx and .pptx formats.

Author: Hackathon Team
Date: November 2025
"""

import io
from typing import Optional
from docx import Document
from pptx import Presentation
from PyPDF2 import PdfReader


class DocumentExtractor:
    """Extract text from various document formats"""

    @staticmethod
    def extract_from_pdf(file_bytes: bytes) -> str:
        """
        Extract text from PDF file

        Args:
            file_bytes: PDF file content as bytes

        Returns:
            Extracted text as string

        Raises:
            ValueError: If PDF is encrypted or cannot be read
        """
        try:
            pdf_file = io.BytesIO(file_bytes)
            reader = PdfReader(pdf_file)

            if reader.is_encrypted:
                raise ValueError("Encrypted PDFs are not supported. Please provide an unencrypted PDF.")

            text_content = []
            for page_num, page in enumerate(reader.pages, 1):
                try:
                    text = page.extract_text()
                    if text.strip():
                        text_content.append(text)
                except Exception as e:
                    print(f"Warning: Could not extract text from page {page_num}: {str(e)}")
                    continue

            full_text = "\n\n".join(text_content)

            if not full_text.strip() or len(full_text.strip()) < 50:
                raise ValueError(
                    "Could not extract sufficient text from PDF. "
                    "The PDF might be scanned/image-based or empty."
                )

            return full_text.strip()

        except Exception as e:
            if isinstance(e, ValueError):
                raise
            raise ValueError(f"Failed to process PDF: {str(e)}")

    @staticmethod
    def extract_from_docx(file_bytes: bytes) -> str:
        """
        Extract text from Word DOCX file

        Args:
            file_bytes: DOCX file content as bytes

        Returns:
            Extracted text as string

        Raises:
            ValueError: If DOCX cannot be read or is empty
        """
        try:
            docx_file = io.BytesIO(file_bytes)
            doc = Document(docx_file)

            text_content = []

            # Extract text from paragraphs
            for paragraph in doc.paragraphs:
                text = paragraph.text.strip()
                if text:
                    text_content.append(text)

            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        cell_text = cell.text.strip()
                        if cell_text:
                            row_text.append(cell_text)
                    if row_text:
                        text_content.append(" | ".join(row_text))

            full_text = "\n\n".join(text_content)

            if not full_text.strip() or len(full_text.strip()) < 50:
                raise ValueError(
                    "Could not extract sufficient text from Word document. "
                    "The document might be empty or contain only images."
                )

            return full_text.strip()

        except Exception as e:
            if isinstance(e, ValueError):
                raise
            raise ValueError(f"Failed to process Word document: {str(e)}")

    @staticmethod
    def extract_from_pptx(file_bytes: bytes) -> str:
        """
        Extract text from PowerPoint PPTX file

        Args:
            file_bytes: PPTX file content as bytes

        Returns:
            Extracted text as string

        Raises:
            ValueError: If PPTX cannot be read or is empty
        """
        try:
            pptx_file = io.BytesIO(file_bytes)
            prs = Presentation(pptx_file)

            text_content = []

            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []

                # Extract title
                if slide.shapes.title:
                    title = slide.shapes.title.text.strip()
                    if title:
                        slide_text.append(f"=== Slide {slide_num}: {title} ===")

                # Extract text from all shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text = shape.text.strip()
                        if text and text != slide.shapes.title.text if slide.shapes.title else True:
                            slide_text.append(text)

                    # Extract text from tables
                    if shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            row_text = []
                            for cell in row.cells:
                                cell_text = cell.text.strip()
                                if cell_text:
                                    row_text.append(cell_text)
                            if row_text:
                                slide_text.append(" | ".join(row_text))

                if slide_text:
                    text_content.append("\n".join(slide_text))

            full_text = "\n\n".join(text_content)

            if not full_text.strip() or len(full_text.strip()) < 50:
                raise ValueError(
                    "Could not extract sufficient text from PowerPoint. "
                    "The presentation might be empty or contain only images."
                )

            return full_text.strip()

        except Exception as e:
            if isinstance(e, ValueError):
                raise
            raise ValueError(f"Failed to process PowerPoint presentation: {str(e)}")

    @staticmethod
    def extract_from_txt(file_bytes: bytes) -> str:
        """
        Extract text from plain text file

        Args:
            file_bytes: TXT file content as bytes

        Returns:
            Extracted text as string

        Raises:
            ValueError: If text file is empty or cannot be decoded
        """
        try:
            # Try UTF-8 first
            try:
                text = file_bytes.decode('utf-8')
            except UnicodeDecodeError:
                # Fallback to latin-1
                text = file_bytes.decode('latin-1')

            text = text.strip()

            if not text or len(text) < 50:
                raise ValueError("Text file is empty or too short (minimum 50 characters required)")

            return text

        except Exception as e:
            if isinstance(e, ValueError):
                raise
            raise ValueError(f"Failed to process text file: {str(e)}")

    @staticmethod
    def extract_text(file_bytes: bytes, filename: str) -> tuple[str, str]:
        """
        Extract text from any supported document format

        Args:
            file_bytes: File content as bytes
            filename: Original filename with extension

        Returns:
            Tuple of (extracted_text, file_type)

        Raises:
            ValueError: If file format is unsupported or extraction fails
        """
        file_extension = filename.lower().split('.')[-1] if '.' in filename else ''

        extractors = {
            'pdf': (DocumentExtractor.extract_from_pdf, 'PDF'),
            'doc': (DocumentExtractor.extract_from_docx, 'Word Document'),
            'docx': (DocumentExtractor.extract_from_docx, 'Word Document'),
            'ppt': (DocumentExtractor.extract_from_pptx, 'PowerPoint'),
            'pptx': (DocumentExtractor.extract_from_pptx, 'PowerPoint'),
            'txt': (DocumentExtractor.extract_from_txt, 'Text File'),
        }

        if file_extension not in extractors:
            supported = ', '.join(extractors.keys())
            raise ValueError(
                f"Unsupported file format: .{file_extension}. "
                f"Supported formats: {supported}"
            )

        extractor_func, file_type = extractors[file_extension]
        text = extractor_func(file_bytes)

        return text, file_type
